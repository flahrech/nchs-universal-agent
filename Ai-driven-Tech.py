import os
from dotenv import load_dotenv
from crewai import Agent, Task, Crew, Process
from crewai_tools import TavilySearchTool
from pptx import Presentation

# Load secure NCHS credentials
load_dotenv()

# Initialize the "Eyes" of the department
search_tool = TavilySearchTool()

# --- THE EXECUTIVE FORMATTER TOOL ---
def create_branded_report(ai_content):
    """Dynamically creates slides based on the AI's section headers."""
    try:
        prs = Presentation('NCHS_Template.pptx')
    except:
        prs = Presentation()

    # Split the AI content by major headers (assuming Markdown headers like '##')
    sections = str(ai_content).split('##')

    for section in sections:
        if len(section).strip() > 10:  # Only process sections with actual content
            # Add a 'Title and Content' slide (usually layout index 1)
            slide_layout = prs.slide_layouts[1] 
            slide = prs.slides.add_slide(slide_layout)
            
            # Extract the first line as the slide title, the rest as body
            lines = section.strip().split('\n')
            slide.shapes.title.text = lines[0] if lines[0] else "NCHS Analysis"
            
            # Put the remaining text into the content placeholder
            body_text = "\n".join(lines[1:])
            slide.placeholders[1].text = body_text

    save_path = 'NCHS_Executive_Briefing.pptx'
    prs.save(save_path)
    return f"Report saved to {save_path}"

# --- AGENT DEFS: The Digital Department ---

researcher = Agent(
    role='Lead Process Researcher',
    goal='Identify cutting-edge AI and Lean trends in top-tier healthcare systems',
    backstory='Senior analyst at NCHS. You find the "What" and "Where" of innovation.',
    tools=[search_tool],
    verbose=True,
    allow_delegation=False
)

financial_analyst = Agent(
    role='Healthcare ROI Specialist',
    goal='Calculate the cost-benefit and ROI of proposed process changes',
    backstory='Former hospital CFO. You ensure every NCHS project is financially sound.',
    verbose=True,
    allow_delegation=True # Can ask the researcher for more data
)

compliance_officer = Agent(
    role='Healthcare Risk & Compliance Lead',
    goal='Ensure all AI and process innovations meet HIPAA and CMS guidelines',
    backstory='Legal expert at NCHS. You protect the organization from risk.',
    verbose=True
)

writer = Agent(
    role='Director of Communications',
    goal='Synthesize research, finance, and compliance into a board-ready memo',
    backstory='The bridge between technical data and C-suite decision-making.',
    verbose=True
)

# --- TASK DEFS: The Workflow ---

research_task = Task(
    description='Analyze AI-driven patient triage impacts on ER wait times in 2025-2026.',
    expected_output='Detailed summary of 3 successful hospital implementations.',
    agent=researcher
)

finance_task = Task(
    description='Based on the research, estimate the potential cost savings for NCHS.',
    expected_output='A brief ROI table or bulleted financial breakdown.',
    agent=financial_analyst
)

compliance_task = Task(
    description='Evaluate the triage AI for HIPAA and patient data privacy risks.',
    expected_output='A "Risk & Mitigation" summary.',
    agent=compliance_officer
)

writing_task = Task(
    description='Combine all findings into a final executive memo for the board.',
    expected_output='A professional 4-section memo in Markdown.',
    agent=writer
)

# --- THE ENGINE: Sequential Process ---

nchs_department = Crew(
    agents=[researcher, financial_analyst, compliance_officer, writer],
    tasks=[research_task, finance_task, compliance_task, writing_task],
    process=Process.sequential # Lean workflow: One stage must pass quality gates to proceed
)

if __name__ == "__main__":
    print("### NCHS Digital Department: Commencing Project Analysis ###")
    
    # 1. Execute the multi-agent workflow
    result = nchs_department.kickoff()
    
    # 2. Export to NCHS Branded PowerPoint
    create_branded_report(str(result))
    
    # 3. Save a local Markdown copy for documentation
    with open("Latest_NCHS_Report.md", "w") as f:
        f.write(str(result))
        
    print("\n\n####################################")
    print("## STRATEGIC ANALYSIS COMPLETE ##")
    print("####################################\n")
    print(f"Deliverables generated: NCHS_Executive_Briefing.pptx and Latest_NCHS_Report.md")